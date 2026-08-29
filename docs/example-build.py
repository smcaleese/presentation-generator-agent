# Example of the python-pptx program the model passes to the `createSlides` tool.
#
# The model is told: "output a COMPLETE Python program that uses python-pptx to
# build the presentation and saves it as 'deck.pptx' in the current working
# directory. Do not read or write any path other than 'deck.pptx'. Output only
# the code."
#
# This is a representative result for the prompt:
#   "A 3-slide intro to what a vector database is, for engineers.
#    Title slide plus two content slides."

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---- theme ----------------------------------------------------------------
BG = RGBColor(0x14, 0x14, 0x24)      # near-black navy
FG = RGBColor(0xF5, 0xF5, 0xF7)      # off-white
MUTED = RGBColor(0x9A, 0x9A, 0xB0)   # grey for subtitles / footer
ACCENT = RGBColor(0x6C, 0x5C, 0xE7)  # periwinkle underline / bullets
FONT = "Arial"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
BLANK = prs.slide_layouts[6]


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    return slide


def textbox(slide, left, top, width, height, lines, *, size, color, bold=False,
            align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=6):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def accent_bar(slide, left, top, width=Inches(1.6)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    bar.shadow.inherit = False


def bullets(slide, items):
    box = slide.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(11.3), Inches(4.4))
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(16)
        dot = p.add_run()
        dot.text = "•  "
        dot.font.name = FONT
        dot.font.size = Pt(22)
        dot.font.color.rgb = ACCENT
        run = p.add_run()
        run.text = item
        run.font.name = FONT
        run.font.size = Pt(22)
        run.font.color.rgb = FG


def content_slide(title, items):
    s = add_slide()
    textbox(s, Inches(1.0), Inches(0.8), Inches(11.3), Inches(1.2),
            [title], size=40, color=FG, bold=True)
    accent_bar(s, Inches(1.0), Inches(1.95))
    bullets(s, items)
    return s


# ---- slide 1 — title ----------------------------------------------------
s1 = add_slide()
textbox(s1, Inches(1.0), Inches(2.6), Inches(11.3), Inches(1.6),
        ["Vector Databases"], size=60, color=FG, bold=True)
accent_bar(s1, Inches(1.05), Inches(3.9))
textbox(s1, Inches(1.0), Inches(4.1), Inches(11.3), Inches(1.0),
        ["An engineering introduction to semantic search at scale"],
        size=26, color=MUTED)
textbox(s1, Inches(1.0), Inches(6.6), Inches(11.3), Inches(0.5),
        ["Embeddings · Similarity search · ANN indexing"],
        size=16, color=MUTED)

# ---- slide 2 ----------------------------------------------------------------
content_slide("What is a vector database?", [
    "Stores high-dimensional embeddings — arrays of floats produced by an "
    "ML model that place semantically similar items near each other.",
    "Primary query is nearest-neighbour: 'given this vector, return the k "
    "closest', not exact key lookup or range scans.",
    "Uses approximate nearest-neighbour (ANN) indexes — HNSW, IVF, PQ — to "
    "trade a little recall for sub-linear query time.",
    "Adds the database parts around that: metadata filtering, upserts, "
    "persistence, sharding, replication.",
])

# ---- slide 3 --------------------------------------------------------------
content_slide("When you'd reach for one", [
    "Semantic search over documents, code, or support tickets.",
    "Retrieval-augmented generation — fetch context for an LLM prompt.",
    "Recommendations and de-duplication by similarity.",
    "Not needed when exact filters or full-text search already answer the "
    "query — a Postgres GIN index or pgvector may be enough.",
])

prs.save("deck.pptx")
print("wrote deck.pptx —", len(prs.slides), "slides")
